import io
import logging
import asyncio
from typing import Optional, List, Dict, Any
from datetime import datetime
import json

from app.core.config import settings
from app.services.storage import minio_service
from app.services.milvus_service import milvus_service

logger = logging.getLogger(__name__)


class DocumentParserService:
    def __init__(self):
        pass
    
    async def parse_and_store(
        self,
        project_id: str,
        document_id: str,
        file_name: str,
        file_data: bytes,
        doc_type: str
    ) -> Dict[str, Any]:
        logger.info(f"Starting to parse and store document: {file_name}")
        
        try:
            parsed_content = await self.parse_document(
                file_data=file_data,
                file_name=file_name
            )
            
            chunks = self._split_content(parsed_content)
            
            if chunks:
                vector_count = await milvus_service.insert_vectors(
                    project_id=project_id,
                    document_id=document_id,
                    document_name=file_name,
                    chunks=chunks,
                    content_type=doc_type
                )
                
                logger.info(f"Stored {vector_count} vectors for document {file_name}")
            
            return {
                "success": True,
                "content": parsed_content,
                "chunks": chunks,
                "metadata": {
                    "file_name": file_name,
                    "doc_type": doc_type,
                    "parsed_at": datetime.utcnow().isoformat(),
                    "chunk_count": len(chunks)
                }
            }
            
        except Exception as e:
            logger.error(f"Error parsing document {file_name}: {e}")
            return {
                "success": False,
                "error": str(e),
                "content": "",
                "chunks": []
            }
    
    async def parse_document(
        self,
        file_data: bytes,
        file_name: str
    ) -> str:
        ext = file_name.lower().split(".")[-1] if "." in file_name else ""
        
        try:
            if ext == "pdf":
                return await self._parse_pdf(file_data)
            elif ext in ["docx", "doc"]:
                return await self._parse_docx(file_data)
            elif ext == "md":
                return file_data.decode("utf-8", errors="ignore")
            elif ext == "txt":
                return file_data.decode("utf-8", errors="ignore")
            elif ext in ["xlsx", "xls"]:
                return await self._parse_excel(file_data)
            elif ext in ["pptx", "ppt"]:
                return await self._parse_pptx(file_data)
            else:
                try:
                    return file_data.decode("utf-8", errors="ignore")
                except:
                    return ""
        except Exception as e:
            logger.error(f"Error in parse_document: {e}")
            return ""
    
    async def _parse_pdf(self, file_data: bytes) -> str:
        try:
            from pypdf import PdfReader
            
            reader = PdfReader(io.BytesIO(file_data))
            text_parts = []
            
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"PDF parse error: {e}")
            return ""
    
    async def _parse_docx(self, file_data: bytes) -> str:
        try:
            from docx import Document
            
            doc = Document(io.BytesIO(file_data))
            text_parts = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                text_parts.append("\n".join(table_text))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"DOCX parse error: {e}")
            return ""
    
    async def _parse_excel(self, file_data: bytes) -> str:
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(io.BytesIO(file_data))
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"## Sheet: {sheet_name}\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell else "" for cell in row]
                    if any(row_text):
                        text_parts.append(" | ".join(row_text))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"Excel parse error: {e}")
            return ""
    
    async def _parse_pptx(self, file_data: bytes) -> str:
        try:
            from pptx import Presentation
            
            prs = Presentation(io.BytesIO(file_data))
            text_parts = []
            
            for i, slide in enumerate(prs.slides):
                text_parts.append(f"## Slide {i+1}\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.error(f"PPTX parse error: {e}")
            return ""
    
    def _split_content(
        self,
        content: str,
        chunk_size: int = 500,
        overlap: int = 50
    ) -> List[Dict[str, Any]]:
        if not content:
            return []
        
        chunks = []
        chunk_num = 0
        
        paragraphs = content.split('\n\n')
        current_chunk = ""
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
                
            if len(current_chunk) + len(para) > chunk_size:
                if current_chunk:
                    chunks.append({
                        "chunk_num": chunk_num,
                        "content": current_chunk.strip(),
                        "char_count": len(current_chunk.strip()),
                        "source": "document"
                    })
                    chunk_num += 1
                current_chunk = para
            else:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para
        
        if current_chunk:
            chunks.append({
                "chunk_num": chunk_num,
                "content": current_chunk.strip(),
                "char_count": len(current_chunk.strip()),
                "source": "document"
            })
        
        return chunks
    
    def get_supported_formats(self) -> List[str]:
        return [
            ".pdf", ".docx", ".doc", ".pptx", ".ppt",
            ".xlsx", ".xls", ".md", ".txt", ".html"
        ]


document_parser = DocumentParserService()
